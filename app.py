# from flask import Flask, request, jsonify
# from flask_cors import CORS, cross_origin
# from pptx import Presentation
# import collections.abc
# from firebase_admin import credentials, initialize_app, storage

# cred = credentials.Certificate("firebase_credentials.json")
# initialize_app(cred, {"storageBucket": "presento-1d9cd.appspot.com"})

# app = Flask(__name__)
# cors = CORS(app)

# @app.route("/", methods=["POST", "GET", "PUT", "OPTIONS"])
# def make():
#     def makePPT(data):
#         def _add_leveled_bullet(_placeholder, _text, level=0):
#             _prg = _placeholder.text_frame.add_paragraph()
#             _prg.level = level
#             _prg.text = _text

#         # Create a presentation object
#         prs = Presentation()
#         title_slide_layout = prs.slide_layouts[0]
#         slide = prs.slides.add_slide(title_slide_layout)
#         title = slide.shapes.title
#         title.text = data["presentationTitle"]

#         for i in range(len(data["slide"])):
#             slide_layout = prs.slide_layouts[1]
#             slide = prs.slides.add_slide(slide_layout)
#             title = slide.shapes.title
#             subtitle = slide.placeholders[1]
#             title.text = data["slide"][i]["title"]
#             for j in range(len(data["slide"][i]["points"])):
#                 _add_leveled_bullet(subtitle, data["slide"][i]["points"][j], 0)
#         # Save the presentation to a file
#         prs.save("my_presentation.pptx")

#     json_data = request.get_json()
#     print(json_data)
#     makePPT(json_data)
#     file_path = "my_presentation.pptx"

#     bucket = storage.bucket()
#     blob = bucket.blob(file_path)
#     blob.upload_from_filename(file_path)
#     blob.make_public()
#     return blob.public_url

# if __name__ == "__main__":
#     app.run(host="0.0.0.0", port=5000, debug=True)
from flask import Flask, send_file, request
from flask_cors import CORS, cross_origin
from pptx import Presentation
import collections.abc

app = Flask(__name__)
cors = CORS(app)
app.config['CORS_HEADERS'] = 'Content-Type'

# # Create a presentation object
# presentation = Presentation()

# # Add a slide with a title and content
# slide_layout = presentation.slide_layouts[
#   1]  # Choose a slide layout (e.g., Title and Content)
# slide = presentation.slides.add_slide(slide_layout)

# title = slide.shapes.title
# title.text = "My First Slide"

# content = slide.placeholders[1]
# content.text = "Hello, PowerPoint!"

# # Save the presentation to a file
# presentation.save("my_presentation.pptx")

def makePPT(data):
  def _add_leveled_bullet(_placeholder, _text, level=0):
    _prg = _placeholder.text_frame.add_paragraph()
    _prg.level = level
    _prg.text = _text
  
  # Create a presentation object
  prs = Presentation()
  title_slide_layout = prs.slide_layouts[0]
  slide = prs.slides.add_slide(title_slide_layout)
  title = slide.shapes.title
  title.text = data['presentationTitle']
  
  for i in range(len(data['slide'])):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = data["slide"][i]["title"]
    for j in range(len(data["slide"][i]["points"])):
      _add_leveled_bullet(subtitle,data["slide"][i]["points"][j] , 0)
  # Save the presentation to a file
  prs.save("SID_presentation.pptx")



@app.route('/',methods=['POST'])
@cross_origin()
def index():
  json_data = request.get_json()
  # return json_data

  makePPT(json_data)
  file_path = 'SID_presentation.pptx'

  return send_file(file_path, as_attachment=True)

# @app.route('/display')
# @cross_origin()
# def index():
#   file_path = 'my_presentation.pptx'
  
#   return send_file(file_path, as_attachment=False)


app.run(host='0.0.0.0', port=5000, debug=True)
